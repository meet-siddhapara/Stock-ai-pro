import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle

def add_content_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

def add_image_slide(prs, title, image_path, summary):
    slide_layout = prs.slide_layouts[5] # Title only layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add image
    if os.path.exists(image_path):
        # Calculate image dimensions to fit nicely
        img = slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    else:
        # Add a placeholder text box if image doesn't exist
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.text = f"[Image missing: {os.path.basename(image_path)}]"
    
    # Add summary below
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = summary
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    add_title_slide(prs, 
                   "Stock Market Screener and Prediction", 
                   "Using Machine Learning and Streamlit\nBy Meet & Team")
                   
    # 2. Objective
    add_content_slide(prs, "Project Objective", [
        "Build an efficient Machine Learning model to predict future values of financial stocks.",
        "Assist retail investors using a third-party web app developed with Streamlit.",
        "Provide critical features including fundamental information, technical indicators, screeners, and pattern recognition.",
        "Integrate historical data and visualize it dynamically."
    ])
    
    # 3. Tech Stack
    add_content_slide(prs, "Technology Stack", [
        "Programming Language: Python",
        "Machine Learning Packages: Scikit-learn, Pandas, Numpy, Tensorflow, Keras",
        "Data Visualization: Matplotlib, Plotly",
        "Technical Indicators: FinTA, TALib",
        "Data Retrieval: yfinance (OHLC data)",
        "Web App Framework: Streamlit"
    ])
    
    # 4. Model Comparison
    add_content_slide(prs, "Model Comparison", [
        "Models Trained: Moving Average, K-Nearest Neighbors, Linear Regression, LSTM.",
        "Tested on: Large Cap (TCS), Mid Cap (Tata Motors), Small Cap (Trident).",
        "Evaluation Metric: Root Mean Square Error (RMSE).",
        "Result: LSTM performed incredibly well with the lowest error rates across all capitalization tiers.",
        "Conclusion: Deployed LSTM as the primary prediction module in the final Web App."
    ])
    
    # 5. Snapshots & Summaries
    
    base_path = r"c:\Users\Meet\Downloads\final year\Stock-market-prediction-and-screener\media"
    
    # Homepage / Web App Overview
    add_image_slide(prs, "Web Application Overview",
                    os.path.join(base_path, "Screenshot (189).png"),
                    "Summary: The interactive dashboard home interface showing real-time selection of stock symbols.")

    # Fundamental Information
    add_image_slide(prs, "Fundamental Information Module",
                    os.path.join(base_path, "Screenshot (195).png"),
                    "Summary: Detailed historical prices, balance sheet, financials, and company fundamentals visually presented.")
                    
    # Technical Indicators
    add_image_slide(prs, "Technical Indicators Analysis",
                    os.path.join(base_path, "Screenshot (205).png"),
                    "Summary: Interactive charts plotting crucial technical indicators like Moving Averages, RSI, MACD, and Bollinger Bands.")
                    
    # Screener
    add_image_slide(prs, "Technical Screener",
                    os.path.join(base_path, "Screenshot (210).png"),
                    "Summary: Filtering parameters like 52-week highs/lows and RSI to strategically screen and select winning stocks.")
                    
    # Pattern Recognition
    add_image_slide(prs, "Pattern Recognition System",
                    os.path.join(base_path, "Screenshot (215).png"),
                    "Summary: Automated detection of Japanese candlestick patterns, offering bullish or bearish signals instantly.")
                    
    # Next-Day Forecasting (LSTM)
    add_image_slide(prs, "Next-Day Price Forecasting",
                    os.path.join(base_path, "tcslstm.png"),
                    "Summary: Displaying our LSTM Deep Learning model's actual vs predicted prices on large tech stocks like TCS.")
                    
    # 6. Limitations & Future Scope
    add_content_slide(prs, "Limitations and Future Scope", [
        "Limitations: Stock market volatility makes extremely precise absolute prediction challenging. Physical, psychological, and irrational market behavior factors persist.",
        "Future Scope: Expansion of prediction horizons. Currently focusing on the 1-day timeframe; future versions can accommodate intra-day (minutes) and weekly/monthly timeframes.",
        "Future Scope: Continual model hyperparameter tuning and integration of natural language sentiment analysis."
    ])
    
    # 7. Thank You
    add_title_slide(prs, "Thank You", "Any Questions?")
    
    out_file = "Stock_Market_Prediction_Presentation.pptx"
    prs.save(out_file)
    print(f"Presentation generated successfully: {out_file}")

if __name__ == '__main__':
    create_presentation()
